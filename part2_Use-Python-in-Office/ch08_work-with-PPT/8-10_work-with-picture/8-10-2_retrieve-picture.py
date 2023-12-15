from pptx import Presentation

from pptx.shapes.picture import Picture
from pptx.shapes.placeholder import PlaceholderPicture

ppt = Presentation("./8-10-1.pptx")

for slide in ppt.slides:
    for shape in slide.shapes:
        pic_type = (Picture, PlaceholderPicture)
        if not isinstance(shape,pic_type):
            continue
        content_type = shape.image.content_type
        print(content_type)
        file_type = content_type.split("/")[-1]
        
        if file_type == "jpeg":
            file_type = "jpg"
        
        print(file_type)
        shape_name = shape.name
        print(shape_name)
        shape_name = str(shape_name).replace(" ", "")
        print(shape_name)
        
        save_path = f"./{shape_name}.{file_type}"
        blob = shape.image.blob
        with open(save_path, "wb") as f:
            f.write(blob)