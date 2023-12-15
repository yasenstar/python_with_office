from pptx import Presentation
from pptx.util import Cm

ppt = Presentation()

slide = ppt.slides.add_slide(ppt.slide_layouts[6])

path = "./galaxy.jpg"

top = left = Cm(3)
width = Cm(20); height = Cm(10)

picture = slide.shapes.add_picture(
    path,
    top,
    left,
    width,
    height
)

ppt.save("./8-10-1.pptx")