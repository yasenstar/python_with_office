from pptx import Presentation
ppt = Presentation("./8-4-4.pptx")

sld_list = list(ppt.slides._sldIdLst)

for i in [1,2]:
    slide_el = sld_list[i]
    ppt.slides._sldIdLst.remove(slide_el)

ppt.save("./8-4-5.pptx")