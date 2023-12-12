from pptx import Presentation

ppt = Presentation()

slide_layout = ppt.slide_layouts[1]

slide1 = ppt.slides.add_slide(slide_layout)
slide2 = ppt.slides.add_slide(slide_layout)
slide3 = ppt.slides.add_slide(slide_layout)
slide4 = ppt.slides.add_slide(slide_layout)
slide5 = ppt.slides.add_slide(slide_layout)

print(ppt.slides.index(slide1))

print(ppt.slides.index(slide2))

print(ppt.slides.index(slide5))

# ppt.save("./8-4-3.pptx")