from pptx import Presentation
ppt = Presentation("./8-4-3.pptx")

slide1 = ppt.slides[2]

slide2 = ppt.slides.get(2)

print(type(slide1)) # <class 'pptx.slide.Slide'>
print(type(slide2)) # <class 'NoneType'>

ppt.save("./8-4-4.pptx")