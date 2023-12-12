from pptx import Presentation
ppt = Presentation()

slide_layout = ppt.slide_layouts[10]

print(type(ppt.slides))

slide = ppt.slides.add_slide(slide_layout)

ppt.slides.add_slide(ppt.slide_layouts[1])

print(type(slide))

ppt.save("./8-4-2.pptx")