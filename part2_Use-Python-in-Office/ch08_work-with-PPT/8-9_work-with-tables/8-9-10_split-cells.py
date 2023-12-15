from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

ppt = Presentation("./8-9-9_2.pptx")

for slide in ppt.slides:
    for shape in slide.shapes:
        if shape.shape_type != MSO_SHAPE_TYPE.TABLE:
            continue
        table = shape.table
        c1 = table.cell(0,1)
        c2 = table.cell(2,2)
        print(c1.is_merge_origin)
        print(c2.is_spanned)
        print(table.cell(1,1).is_spanned)
        
        if c1.is_merge_origin:
            c1.split()
        
        print(c1.is_merge_origin)
        print(c2.is_spanned)
        print(table.cell(1,1).is_spanned)

ppt.save("./8-9-10.pptx")