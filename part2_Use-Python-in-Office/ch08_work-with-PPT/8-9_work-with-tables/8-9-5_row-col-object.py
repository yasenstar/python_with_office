from pptx import Presentation
from pptx.util import Cm

ppt = Presentation()

slide = ppt.slides.add_slide(ppt.slide_layouts[6])

row = 4; column = 3

left = top = Cm(3)
width = Cm(20); height = Cm(10)

table_frame = slide.shapes.add_table(
    row,
    column,
    top,
    left,
    width,
    height
)
print(type(table_frame))

table = table_frame.table
print(type(table))

rows = table.rows
columns = table.columns

print(type(rows))
print(type(columns))

row = rows[0]
column = columns[0]

print(type(row))
print(type(column))

row.height = Cm(5)
column.width = Cm(5)

ppt.save("./8-9-5.pptx")