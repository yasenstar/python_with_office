from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Cm, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT, MSO_VERTICAL_ANCHOR

ppt = Presentation()

slide = ppt.slides.add_slide(ppt.slide_layouts[6])

table = slide.shapes.add_table(
    4,
    3,
    Cm(3),
    Cm(3),
    Cm(20),
    Cm(10)
).table

text_frame = table.cell(0,0).text_frame

text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

paragraph = text_frame.paragraphs[0]

paragraph.text = "PYTHON"

paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

run = paragraph.runs[0]

run.font.color.rgb = RGBColor.from_string("FF0000")

run.font.size = Pt(30)

fill = table.cell(0,0).fill

fill.solid()

fill.fore_color.rgb = RGBColor.from_string("708069")



ppt.save("./8-9-8.pptx")