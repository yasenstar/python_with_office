from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

ppt = Presentation("./8-9-2.pptx")

for slide in ppt.slides:
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            print("Current Shape is a Table")
            print(type(shape))