from pptx import Presentation
from pptx.util import Cm

ppt = Presentation()

slide = ppt.slides.add_slide(ppt.slide_layouts[6])

table = slide.shapes.add_table(
    4,
    3,
    Cm(3),
    Cm(3),
    Cm(20),
    Cm(10)
).table

cell = table.cell(1,0)

cell.text = "Python"

print(cell.text)

table.cell(1,1).text = "Office"
table.cell(1,2).text = "PPT"

ppt.save("./8-9-7.pptx")