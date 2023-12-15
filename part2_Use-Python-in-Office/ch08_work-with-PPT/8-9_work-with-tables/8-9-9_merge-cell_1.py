from pptx import Presentation
from pptx.util import Cm

ppt = Presentation()

slide = ppt.slides.add_slide(ppt.slide_layouts[6])

table = slide.shapes.add_table(
    4,
    3,
    Cm(3),
    Cm(3),
    Cm(20),
    Cm(10)
).table

i = 1
for row in table.rows:
    for cell in row.cells:
        cell.text = str(i)
        i += 1

ppt.save("./8-9-9_1.pptx")