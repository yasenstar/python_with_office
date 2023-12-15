from pptx import Presentation
from pptx.util import Cm

ppt = Presentation()

slide = ppt.slides.add_slide(ppt.slide_layouts[6])

table = slide.shapes.add_table(
    4,
    3,
    Cm(3),
    Cm(3),
    Cm(20),
    Cm(10)
).table

cell = table.cell(0,0)

print(type(cell))

cells = table.rows[0].cells

print(type(cells))

for cell in cells:
    print(type(cell))
    cell.text = "test"

ppt.save("./8-9-6.pptx")