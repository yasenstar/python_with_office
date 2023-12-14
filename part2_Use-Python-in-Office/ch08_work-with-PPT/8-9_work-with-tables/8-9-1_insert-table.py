from pptx import Presentation
from pptx.util import Cm

ppt = Presentation()

slide = ppt.slides.add_slide(ppt.slide_layouts[6])

row = 4; column = 3

left = top = Cm(3)
width = Cm(20); height = Cm(10)

table_frame = slide.shapes.add_table(
    row,
    column,
    top,
    left,
    width,
    height
)
print(type(table_frame))

ppt.save("./8-9-1.pptx")