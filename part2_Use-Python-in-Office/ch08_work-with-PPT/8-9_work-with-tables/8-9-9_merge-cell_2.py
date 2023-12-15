from pptx import Presentation
from pptx.util import Cm

ppt = Presentation()

slide = ppt.slides.add_slide(ppt.slide_layouts[6])

table = slide.shapes.add_table(
    4,
    3,
    Cm(3),
    Cm(3),
    Cm(20),
    Cm(10)
).table

i = 1
for row in table.rows:
    for cell in row.cells:
        cell.text = str(i)
        i += 1

c1 = table.cell(0,1)
c2 = table.cell(2,2)

print(c1.is_merge_origin)
print(c2.is_spanned)
print(table.cell(1,1).is_spanned)

c1.merge(c2)

print(c1.is_merge_origin)
print(c2.is_spanned)
print(table.cell(1,1).is_spanned)


ppt.save("./8-9-9_2.pptx")
