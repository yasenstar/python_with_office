from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER_TYPE

ppt = Presentation('./test.pptx')

slide = ppt.slides.add_slide(ppt.slide_layouts[3])

for placeholder in slide.placeholders:
    phf = placeholder.placeholder_format
    if phf.type != PP_PLACEHOLDER_TYPE.TABLE:
        continue
    table_frame = placeholder.insert_table(
        rows = 10,
        cols = 5
    )

ppt.save("./8-9-3.pptx")