from pptx import Presentation
ppt = Presentation()

layout = ppt.slide_layouts[4]

slide = ppt.slides.add_slide(layout)

shapes = slide.shapes

print(type(shapes))

for shape in shapes:
    print(shape.name)
    
    # Title 1
    # Text Placeholder 2
    # Content Placeholder 3
    # Text Placeholder 4
    # Content Placeholder 5

print(len(shapes))

shape1 = shapes[1]
shape2 = shapes[4]

print(type(shape1), shape1.name)
print(type(shape2), shape2.name)

# ppt.save("./8-5-2.pptx")