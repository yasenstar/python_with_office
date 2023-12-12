from pptx import Presentation
from pptx.util import Cm
from pptx.dml.color import RGBColor, ColorFormat

ppt = Presentation()

slide = ppt.slides.add_slide(ppt.slide_layouts[4])

left = top = width = height = Cm(5)

shape = slide.shapes.add_shape(5, left, top, width, height)

fill = shape.fill

print(type(fill))

fill.solid()

fill.fore_color.rgb = RGBColor.from_string("F0FFF0")

fill.patterned()

print(fill.back_color.rgb)

fill.background()

ppt.save("./8-5-4.pptx")