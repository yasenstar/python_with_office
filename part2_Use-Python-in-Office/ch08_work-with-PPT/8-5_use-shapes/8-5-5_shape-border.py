from pptx import Presentation
from pptx.util import Cm
from pptx.dml.color import RGBColor

ppt = Presentation()

slide = ppt.slides.add_slide(ppt.slide_layouts[1])

left = top = width = height = Cm(6)

shape = slide.shapes.add_shape(10, left, top, width, height)

line = shape.line

print(type(line))

line.color.rgb = RGBColor.from_string("00FF00")

line.color.brightness = -0.5

line.width = Cm(1)

ppt.save("./8-5-5.pptx")