from pptx import Presentation
from pptx.util import Cm
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

ppt = Presentation()

slide = ppt.slides.add_slide(ppt.slide_layouts[0])

autoshape_type_id = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE

left = top = Cm(4.0)

width = height = Cm(5.0)

shape = slide.shapes.add_shape(autoshape_type_id, left, top, width, height)

print(type(shape), shape.name)

ppt.save("./8-5-3.pptx")