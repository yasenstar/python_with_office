from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import PP_PLACEHOLDER_TYPE

ppt = Presentation("./test.pptx")

slide = ppt.slides.add_slide(ppt.slide_layouts[2])

for placeholder in slide.placeholders:
    phf = placeholder.placeholder_format
    # print(phf.type)
    if phf.type != PP_PLACEHOLDER_TYPE.CHART:
        continue
    chart_data = ChartData()
    chart_data.categories = ['Q1', 'Q2', 'Q3', 'Q4']
    chart_data.add_series('Digital', (36.6, 21.1, 15.9, 20.4))
    graphic_frame = placeholder.insert_chart(
        XL_CHART_TYPE.PIE,
        chart_data
    )
    graphic_frame.chart.has_legend = True

ppt.save("./8-8-5.pptx")