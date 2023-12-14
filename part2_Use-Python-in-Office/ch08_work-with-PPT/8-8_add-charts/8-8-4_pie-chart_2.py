from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_DATA_LABEL_POSITION
from pptx.util import Cm

chart_data = ChartData()
chart_data.categories = ['Q1', 'Q2', 'Q3', 'Q4']
sale_volume = [36.6, 21.1, 15.9, 20.4]
percent = [volume / sum(sale_volume) for volume in sale_volume]
# print(percent)

chart_data.add_series('Digital', percent)

ppt = Presentation()

slide = ppt.slides.add_slide(ppt.slide_layouts[6])

x = y = Cm(3)
width = Cm(20); height = Cm(10)

graphic_frame = slide.shapes.add_chart(
    XL_CHART_TYPE.PIE,
    x,
    y,
    width,
    height,
    chart_data
)

chart = graphic_frame.chart
chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.RIGHT
chart.plots[0].has_data_labels = True
data_labels = chart.plots[0].data_labels
data_labels.number_format = '0%'
data_labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END

ppt.save("./8-8-4_2.pptx")