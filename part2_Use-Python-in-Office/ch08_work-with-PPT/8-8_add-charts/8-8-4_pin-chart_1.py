from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Cm

chart_data = ChartData()
chart_data.categories = ['Q1', 'Q2', 'Q3', 'Q4']
sale_volume = [36.6, 21.1, 15.9, 20.4]

chart_data.add_series('Digital', sale_volume)

ppt = Presentation()

slide = ppt.slides.add_slide(ppt.slide_layouts[6])

x = y = Cm(3)
width = Cm(20); height = Cm(10)

chart = slide.shapes.add_chart(
    XL_CHART_TYPE.PIE,
    x,
    y,
    width,
    height,
    chart_data
)

ppt.save("./8-8-4_1.pptx")