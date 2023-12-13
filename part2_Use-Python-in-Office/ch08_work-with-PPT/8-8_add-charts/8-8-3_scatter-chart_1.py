from pptx import Presentation
from pptx.chart.data import XyChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Cm

ppt = Presentation()

chart_data = XyChartData()

sale_info = [
    {
        "category": "digital",
        "sale_volume": [36.6, 21.1, 15.9, 20.4]
    },
    {
        "category": "life",
        "sale_volume": [65.5, 52.1, 22.3, 35.3]
    },
    {
        "category": "study",
        "sale_volume": [10.0, 3.1, 9.8, 3.2]
    },
]

for info in sale_info:
    series = chart_data.add_series(info.get("category"))
    for quarter, volume in enumerate(info.get("sale_volume")):
        series.add_data_point(quarter + 1, volume, number_format='"Quarter "0')

slide = ppt.slides.add_slide(ppt.slide_layouts[6])

x = y = Cm(3)
width = Cm(20); height = Cm(10)

chart_frame = slide.shapes.add_chart(
    XL_CHART_TYPE.XY_SCATTER,
    x,
    y,
    width,
    height,
    chart_data
)

chart = chart_frame.chart
x_axis = chart.category_axis
x_axis.has_title = True
x_axis.axis_title.text_frame.text = "Quarter"
y_axis = chart.value_axis
y_axis_has_title = True
y_axis.axis_title.text_frame.text = "Sales Volume"
chart.has_legend = True # by default is True already

ppt.save("./8-8-3_1.pptx")