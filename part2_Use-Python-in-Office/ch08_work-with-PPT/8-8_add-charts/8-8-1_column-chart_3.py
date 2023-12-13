from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_TICK_MARK, XL_LEGEND_POSITION
from pptx.util import Cm

ppt = Presentation()

chart_data = CategoryChartData()

chart_data.categories = ['Digital Entertainment', 'Life Styles', 'Studying']

chart_data.add_series('Q1', (36.6, 65.5, 10.0))
chart_data.add_series('Q2', (21.1, 52.1, 3.1))
chart_data.add_series('Q3', (15.9, 22.3, 9.8))
chart_data.add_series('Q4', (20.4, 35.3, 3.2))

slide = ppt.slides.add_slide(ppt.slide_layouts[6])

x = y = Cm(3)

width = Cm(20)
height = Cm(10)

graphic_frame = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    x,
    y,
    width,
    height,
    chart_data
)

chart = graphic_frame.chart

category_axis = chart.category_axis

category_axis.has_major_gridlines = True
category_axis.has_minor_gridlines = True
category_axis.tick_labels.font.italic = True
category_axis.tick_labels.font.size = Cm(0.5)

value_axis = chart.value_axis

value_axis.maximum_scale = 100
value_axis.minor_tick_mark = XL_TICK_MARK.INSIDE
value_axis.has_minor_gridlines = True

tick_labels = value_axis.tick_labels
tick_labels.number_format = '0"万"'
tick_labels.font.bold = True
tick_labels.font.size = Cm(0.5)

chart.has_legend = True

chart.legend.position = XL_LEGEND_POSITION.TOP

chart.legend.include_in_layout = False

ppt.save("./8-8-1_3.pptx")