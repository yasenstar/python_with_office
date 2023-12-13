from pptx import Presentation
from pptx.chart.data import XyChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Cm

ppt = Presentation()

chart_data = XyChartData()

series_1 = chart_data.add_series("series one")
series_1.add_data_point(0.7, 2.7)
series_1.add_data_point(1.8, 3.2)
series_1.add_data_point(2.6, 0.8)

series_2 = chart_data.add_series("series two")
series_2.add_data_point(1.3, 3.7)
series_2.add_data_point(2.7, 2.3)
series_2.add_data_point(1.6, 1.8)

slide = ppt.slides.add_slide(ppt.slide_layouts[6])

x = y = Cm(3)
width = Cm(20); height = Cm(10)

slide.shapes.add_chart(
    XL_CHART_TYPE.XY_SCATTER,
    x,
    y,
    width,
    height,
    chart_data
)

ppt.save("./8-8-3.pptx")