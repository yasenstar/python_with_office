from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Cm

ppt = Presentation()

chart_data = CategoryChartData()

chart_data.categories = ['Digital Entertainment', 'Life Styles', 'Studying']

chart_data.add_series('Q1', (36.6, 65.5, 10.0))
chart_data.add_series('Q2', (21.1, 52.1, 3.1))
chart_data.add_series('Q3', (15.9, 22.3, 9.8))
chart_data.add_series('Q4', (20.4, 35.3, 3.2))

slide = ppt.slides.add_slide(ppt.slide_layouts[6])

x = y = Cm(3)

width = Cm(20)
height = Cm(10)

graphic_frame = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    x,
    y,
    width,
    height,
    chart_data
)

ppt.save("./8-8-1_1.pptx")