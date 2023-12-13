from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.util import Cm

ppt = Presentation()

chart_data = CategoryChartData()

chart_data.categories = ['Q1', 'Q2', 'Q3', 'Q4']

chart_data.add_series('Digital', (36.6, 21.1, 15.9, 20.4))
chart_data.add_series('Life', (65.5, 52.1, 22.3, 35.3))
chart_data.add_series('Study', (10.0, 3.1, 9.8, 3.2))

slide = ppt.slides.add_slide(ppt.slide_layouts[6])

x = y = Cm(3)
width = Cm(20); height = Cm(10)

graphic_frame = slide.shapes.add_chart(
    XL_CHART_TYPE.LINE,
    x,
    y,
    width,
    height,
    chart_data
)

chart = graphic_frame.chart
chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.TOP
chart.legend.include_in_layout = False

ppt.save("./8-8-2.pptx")