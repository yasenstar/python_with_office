from pptx import Presentation
from pptx.enum.text import MSO_VERTICAL_ANCHOR
from pptx.util import Cm

ppt = Presentation()

slide = ppt.slides.add_slide(ppt.slide_layouts[0])

shape = slide.shapes[0]

text_frame = shape.text_frame

text_frame.text = "Office遇见了Python"

# text_frame.margin_left = Cm(5)
# text_frame.margin_top = Cm(2)

text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP

text_box = slide.shapes.add_textbox(Cm(3), Cm(5), Cm(10), Cm(5))

text_frame1 = text_box.text_frame
text_frame1.text = "Office遇上了Python，从此它们就幸福的在一起了。Office遇上了Python，从此它们就幸福的在一起了。"
# text_frame1.word_wrap = True

ppt.save("./8-7-6.pptx")