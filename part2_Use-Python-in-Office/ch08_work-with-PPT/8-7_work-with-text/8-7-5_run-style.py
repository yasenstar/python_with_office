from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Cm

ppt = Presentation()

slide = ppt.slides.add_slide(ppt.slide_layouts[0])

shape = slide.shapes[0]

paragraph = shape.text_frame.paragraphs[0]
paragraph.add_run().text = "Office"
paragraph.add_run().text = "遇上了"

run = paragraph.add_run()
run.text = "Python"
run.font.name = "微软雅黑"
run.font.size = Cm(3)
run.font.bold = True
run.font.italic = True
run.font.underline = True
run.font.color.rgb = RGBColor.from_string("0000FF")
run.hyperlink.address = "https://python.org"

for run in paragraph.runs[:-1]:
    run.font.name = "楷体"
    run.font.size = Cm(4)
    run.font.color.rgb = RGBColor.from_string("FF0000")

ppt.save("./8-7-5.pptx")