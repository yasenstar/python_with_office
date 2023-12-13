from pptx import Presentation

ppt = Presentation()

slide = ppt.slides.add_slide(ppt.slide_layouts[0])

shape = slide.shapes[0]

if shape.has_text_frame:
    text_frame = shape.text_frame
    p1 = text_frame.paragraphs[0]
    p1.text = "this is the default paragraph"
    p2 = text_frame.add_paragraph()
    p2.text = "this is the new added paragraph"
    text_frame.add_paragraph().text = "the 3rd paragraph"
    print(len(text_frame.paragraphs))

ppt.save("./8-7-3_2.pptx")