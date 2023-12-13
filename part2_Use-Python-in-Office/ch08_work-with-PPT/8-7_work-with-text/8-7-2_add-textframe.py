from pptx import Presentation
from pptx.util import Cm
ppt = Presentation()

slide = ppt.slides.add_slide(ppt.slide_layouts[0])

left = top = Cm(3)
width = Cm(15)
height = Cm(5)

text_box = slide.shapes.add_textbox(left, top, width, height)

print(type(text_box), text_box.shape_type)

# Result: <class 'pptx.shapes.autoshape.Shape'> TEXT_BOX (17)

text_frame = text_box.text_frame
print(type(text_frame))

# Result: <class 'pptx.text.text.TextFrame'>

ppt.save("./8-7-2.pptx")