from pptx import Presentation
ppt = Presentation()

slide = ppt.slides.add_slide(ppt.slide_layouts[4])

for shape in slide.shapes:
    if not shape.has_text_frame:
        continue
    text_frame = shape.text_frame
    print(type(text_frame))