from pptx import Presentation

ppt = Presentation()

slide = ppt.slides.add_slide(ppt.slide_layouts[7])

for shape in slide.shapes:
    if not shape.has_text_frame:
        continue
    text_frame = shape.text_frame
    paragraph = text_frame.paragraphs[0]
    paragraph.text = "This is my first Text Frame"

ppt.save("./8-7-3.pptx")