from pptx import Presentation
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.util import Cm

ppt = Presentation()

slide = ppt.slides.add_slide(ppt.slide_layouts[0])

shape = slide.shapes[0]

paragraph = shape.text_frame.paragraphs[0]

paragraph.clear()

paragraph.text = "Office meets Python\nOffice遇见Python"

paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT

paragraph.level = 8

# paragraph.line_spacing = 2.0
paragraph.line_spacing = Cm(5)

ppt.save("./8-7-4.pptx")