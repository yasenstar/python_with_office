from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER_TYPE

ppt = Presentation()

slide = ppt.slides.add_slide(ppt.slide_layouts[8])

for placeholder in slide.placeholders:
    phf = placeholder.placeholder_format
    if phf.type == PP_PLACEHOLDER_TYPE.TITLE:
        placeholder.text = "This is A Title"
    elif phf.type == PP_PLACEHOLDER_TYPE.BODY:
        placeholder.text = "This is the normal text box"
    elif phf.type == PP_PLACEHOLDER_TYPE.PICTURE:
        placeholder.insert_picture("./galaxy.jpg")

ppt.save("./8-6-4_1.pptx")