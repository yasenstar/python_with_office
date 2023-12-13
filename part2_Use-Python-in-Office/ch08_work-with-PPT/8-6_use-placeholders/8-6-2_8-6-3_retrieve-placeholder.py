from pptx import Presentation

ppt = Presentation()

slide = ppt.slides.add_slide(ppt.slide_layouts[10])

for shape in slide.shapes:
    print(type(shape), shape.shape_type, shape.is_placeholder)

for placeholder in slide.placeholders:
    print(placeholder.shape_type, placeholder.name)
    phf = placeholder.placeholder_format
    print(phf.idx, phf.type)